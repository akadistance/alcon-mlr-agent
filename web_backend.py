"""
Flask backend for EyeQ web interface
Uses the shared agent runtime (decoupled from main.py)
"""

from flask import Flask, request, jsonify
from pydantic import BaseModel, ValidationError
from flask_cors import CORS
import json
import os
import uuid
from datetime import datetime, timezone
import traceback
from dotenv import load_dotenv
import pdfplumber
from pptx import Presentation

# Import shared agent runtime
from agent_runtime import agent_executor, ConversationState, comprehensive_agent_executor
from tools import read_docx

# Import OCR utilities
from ocr_utils import extract_text_from_image, is_image_file

# Import citation utilities
from regulatory_citations import extract_citations_from_text, add_citations_to_response

# Import document comparison utilities
from document_comparison import (
    calculate_text_similarity,
    extract_claims_from_text,
    compare_claims,
    analyze_compliance_differences,
    format_comparison_report
)

# Import regulatory knowledge base
from regulatory_knowledge_base import (
    search_knowledge_base,
    get_scenario_guidance,
    get_regulatory_topic,
    format_knowledge_base_response,
    COMPLIANCE_SCENARIOS,
    REGULATORY_KNOWLEDGE
)

# Import collaboration features
from collaboration import (
    create_shareable_link,
    add_comment,
    get_shared_analysis,
    create_review_workflow,
    export_conversation,
    get_collaboration_stats
)

# Load environment variables
load_dotenv()

app = Flask(__name__)
CORS(app)

print("Backend initialized successfully!")
print("Using EXACT working agent system from main.py")
print("Data persistence: localStorage (frontend)")

# Configuration
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Global conversation state for web sessions
conversation_sessions = {}

def get_or_create_session(session_id="default"):
    """Get or create a conversation session"""
    if session_id not in conversation_sessions:
        conversation_sessions[session_id] = ConversationState()
    return conversation_sessions[session_id]

def extract_and_parse_citations(text):
    """
    Extract citations from text and parse them into structured format
    
    Args:
        text: Response text containing citations
        
    Returns:
        Dictionary mapping citation numbers to citation details
    """
    import re
    
    citations_dict = {}
    
    # Look for References section at the end
    references_match = re.search(r'References?:\s*\n((?:\[\d+\].*\n?)+)', text, re.IGNORECASE)
    
    if references_match:
        references_text = references_match.group(1)
        
        # Parse individual citations: [1] Title: URL
        citation_pattern = r'\[(\d+)\]\s*([^:]+):\s*(https?://[^\s\n]+)'
        matches = re.findall(citation_pattern, references_text)
        
        for match in matches:
            number, title, url = match
            citations_dict[int(number)] = {
                "number": int(number),
                "title": title.strip(),
                "url": url.strip()
            }
    
    return citations_dict

def format_response_for_chat(response_text, analysis_data):
    """Format the AI response for clean, ChatGPT-style presentation"""
    if not analysis_data:
        return response_text
    
    # Create a clean, structured response like ChatGPT
    formatted_parts = []
    
    # Main analysis header
    formatted_parts.append("## Compliance Analysis Results")
    formatted_parts.append("")
    
    # Quick overview with key findings
    issues_count = len(analysis_data.get('issues', []))
    approved_count = len([c for c in analysis_data.get('approved_claims', []) if c != "No approved claims found."])
    
    formatted_parts.append("### Key Findings:")
    formatted_parts.append(f"• **{approved_count}** approved claims identified")
    formatted_parts.append(f"• **{issues_count}** compliance issues found")
    formatted_parts.append(f"• Disclaimers: {'Present' if analysis_data.get('disclaimer_present') else 'Missing'}")
    formatted_parts.append("")
    
    # Issues section (if any)
    if analysis_data.get('issues') and len(analysis_data['issues']) > 0:
        formatted_parts.append("### 🚨 Compliance Issues:")
        for i, issue in enumerate(analysis_data['issues'], 1):
            issue_type = issue.get('issue', '').replace('_', ' ').title()
            description = issue.get('description', '')
            suggestion = issue.get('suggestion', '')
            
            formatted_parts.append(f"**{i}. {issue_type}**")
            formatted_parts.append(f"   • {description}")
            if suggestion:
                formatted_parts.append(f"   • **Recommendation:** {suggestion}")
            formatted_parts.append("")
    
    # Approved claims (if any)
    if analysis_data.get('approved_claims') and analysis_data['approved_claims'][0] != "No approved claims found.":
        formatted_parts.append("### Approved Claims:")
        for i, claim in enumerate(analysis_data['approved_claims'], 1):
            formatted_parts.append(f"{i}. {claim}")
        formatted_parts.append("")
    
    # Next steps
    formatted_parts.append("### Next Steps:")
    if issues_count > 0:
        formatted_parts.append("• Review and address the compliance issues identified above")
        formatted_parts.append("• Modify content to align with FDA/FTC guidelines")
    if not analysis_data.get('disclaimer_present'):
        formatted_parts.append("• Add appropriate disclaimers and safety information")
    if approved_count > 0:
        formatted_parts.append("• Approved claims can be used as-is in your materials")
    
    formatted_parts.append("")
    formatted_parts.append("*View detailed scoring and metrics in the Accuracy Panel*")
    
    return "\n".join(formatted_parts)

def clean_general_response(response_text):
    """Clean up general AI responses to remove any analysis formatting"""
    if not response_text or not response_text.strip():
        return "I'm here to help with compliance questions. How can I assist you today?"
    
    import re
    
    # Remove any analysis-related patterns
    cleaned_text = response_text
    
    # Remove Summary sections and everything after
    summary_patterns = [
        r'\n\s*Summary\s*\n.*$',
        r'\n\s*Summary:.*$', 
        r'\n\s*## Summary.*$',
        r'\n\s*### Summary.*$'
    ]
    
    for pattern in summary_patterns:
        cleaned_text = re.sub(pattern, '', cleaned_text, flags=re.DOTALL | re.IGNORECASE)
    
    # Remove analysis completion messages
    analysis_patterns = [
        r'\n\s*Analysis completed\s*$',
        r'\n\s*Analysis complete\s*$',
        r'\n\s*Review completed\s*$',
        r'\n\s*Compliance check completed\s*$'
    ]
    
    for pattern in analysis_patterns:
        cleaned_text = re.sub(pattern, '', cleaned_text, flags=re.IGNORECASE)
    
    # Remove any lines that start with analysis indicators
    lines = cleaned_text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        line_stripped = line.strip().lower()
        # Skip lines that look like analysis headers or summaries
        if (line_stripped.startswith('summary') or 
            line_stripped.startswith('analysis') or
            line_stripped.startswith('approved claims') or
            line_stripped.startswith('compliance issues') or
            line_stripped == 'summary' or
            line_stripped == 'analysis completed'):
            continue
        cleaned_lines.append(line)
    
    cleaned_text = '\n'.join(cleaned_lines).strip()
    
    # If the response is empty after cleaning, provide a default
    if not cleaned_text:
        return "I'm here to help with compliance questions. How can I assist you today?"
    
    return cleaned_text

def parse_structured_response(response_text):
    """Parse structured text response into analysis data"""
    import re
    
    # Remove summary sections from response text before parsing
    cleaned_text = clean_general_response(response_text)
    
    analysis_data = {
        "summary": "Analysis completed",
        "approved_claims": [],
        "issues": [],
        "disclaimer_present": False
    }
    
    try:
        # Extract approved claims
        approved_match = re.search(r'Approved claims?:\s*(.*?)(?=\n\n|\nPotential issues?:|\nIssues?:|$)', response_text, re.DOTALL | re.IGNORECASE)
        if approved_match:
            claims_text = approved_match.group(1).strip()
            # Extract bullet points or numbered items
            claims = re.findall(r'[•\-\*]\s*(.+?)(?=\n[•\-\*]|\n\d+\.|\n\n|$)', claims_text, re.DOTALL)
            if not claims:
                # Try numbered format
                claims = re.findall(r'\d+\.\s*(.+?)(?=\n\d+\.|\n\n|$)', claims_text, re.DOTALL)
            if claims:
                analysis_data["approved_claims"] = [claim.strip() for claim in claims]
            elif claims_text and "No approved claims" not in claims_text:
                analysis_data["approved_claims"] = [claims_text.strip()]
        
        # Extract issues
        issues_match = re.search(r'(?:Potential issues?|Issues?):\s*(.*?)(?=\n\n|$)', response_text, re.DOTALL | re.IGNORECASE)
        if issues_match:
            issues_text = issues_match.group(1).strip()
            # Parse numbered issues - handle both formats
            issue_matches = re.findall(r'(\d+)\.\s*([^\s]+(?:\s+[^\s]+)*?)\s+Description:\s*(.+?)(?=\n\d+\.|\n\n|$)', issues_text, re.DOTALL)
            
            if not issue_matches:
                # Fallback to simpler format
                issue_matches = re.findall(r'(\d+)\.\s*([^:]+):\s*(.+?)(?=\n\d+\.|\n\n|$)', issues_text, re.DOTALL)
            
            for match in issue_matches:
                issue_num, issue_type, description = match
                
                # Extract suggestion if present
                suggestion_match = re.search(r'Suggestion:\s*(.+?)(?=Reference:|$)', description, re.DOTALL)
                suggestion = suggestion_match.group(1).strip() if suggestion_match else "Review and revise as needed"
                
                # Extract reference if present
                reference_match = re.search(r'Reference:\s*(.+?)$', description, re.DOTALL)
                reference = reference_match.group(1).strip() if reference_match else "https://www.fda.gov/regulatory-information/search-fda-guidance-documents"
                
                # Clean up description
                clean_description = re.sub(r'\s*Suggestion:.*$', '', description, flags=re.DOTALL).strip()
                clean_description = re.sub(r'\s*Reference:.*$', '', clean_description, flags=re.DOTALL).strip()
                
                analysis_data["issues"].append({
                    "issue": issue_type.strip().lower().replace(' ', '_'),
                    "description": clean_description,
                    "suggestion": suggestion,
                    "reference": reference
                })
        
        # Check for disclaimers
        disclaimer_keywords = ['consult', 'results may vary', 'individual results', 'see your doctor']
        analysis_data["disclaimer_present"] = any(keyword in response_text.lower() for keyword in disclaimer_keywords)
        
        # If no claims found, set default
        if not analysis_data["approved_claims"]:
            analysis_data["approved_claims"] = ["No approved claims found."]
            
    except Exception as e:
        print(f"Error parsing structured response: {e}")
        # Return basic structure on error
        analysis_data = {
            "summary": "Analysis completed",
            "approved_claims": ["No approved claims found."],
            "issues": [],
            "disclaimer_present": False
        }
    
    return analysis_data

def format_agent_response(raw_response):
    """Format agent response for web display - EXACT same logic as main.py"""
    try:
        # EXACT same response handling from main.py
        output = raw_response["output"]
        
        full_text = ""
        if isinstance(output, list):
            for item in output:
                if isinstance(item, dict) and 'text' in item:
                    full_text += item['text'] + "\n"
                elif isinstance(item, str):
                    full_text += item + "\n"
        elif isinstance(output, dict) and 'text' in output:
            full_text = output['text']
        elif isinstance(output, str):
            full_text = output
        else:
            full_text = str(output)
        
        full_text = full_text.strip()
        
        # Skip reasoning lines - EXACT same logic
        clean_text = "\n".join([
            line for line in full_text.splitlines() 
            if not line.strip().startswith(('Intent:', 'Received', 'Classifying', 'Provided', 'Greeting received', 'Multiple greeting'))
        ]).strip()
        
        # Check for JSON and parse if present - EXACT same logic
        json_start = clean_text.find('{')
        json_end = clean_text.rfind('}') + 1
        if json_start != -1 and json_end > json_start:
            try:
                json_str = clean_text[json_start:json_end]
                structured = json.loads(json_str)
                response = structured.get('summary', 'Compliance analysis results:') + "\n\n"
                
                # Approved claims as simple bullets
                if 'approved_claims' in structured and structured['approved_claims']:
                    response += "Approved claims:\n" + "\n".join([f"• {claim.strip()}" for claim in structured['approved_claims'] if claim.strip()]) + "\n\n"
                
                # Issues as numbered with fields, skipping incomplete
                if 'issues' in structured and structured['issues']:
                    valid_issues = [iss for iss in structured['issues'] if iss.get('description') and iss.get('issue')]
                    if valid_issues:
                        response += "Potential issues:\n"
                        for i, issue in enumerate(valid_issues, 1):
                            response += f"{i}. {issue.get('issue', 'Unnamed issue')}\n   Description: {issue.get('description', 'No description')}\n   Suggestion: {issue.get('suggestion', 'No suggestion')}\n   Reference: {issue.get('reference', 'No reference')}\n\n"
                    else:
                        response += "No potential issues identified.\n"
                
                return response
            except json.JSONDecodeError as e:
                print(f"Parsing error: {str(e)}")
                return clean_text  # Fallback
        else:
            return clean_text
            
    except Exception as e:
        print(f"Error formatting response: {e}")
        return "I apologize, but I encountered an error processing your request."

# ===== MAIN API ENDPOINTS =====

@app.route('/api/analyze', methods=['POST']) # runs comprehensive analysis on promotional materials, normal chat for greetings
def analyze():
    """Main analysis endpoint - performs comprehensive MLR analysis on promotional content, normal chat otherwise"""
    try:
        data = request.get_json()
        message = data.get('message', '')
        session_id = data.get('session_id', 'default')
        streaming = data.get('streaming', True)  # Enabled by default for better UX
        
        if not message:
            return jsonify({"error": "No message provided"}), 400
        
        print(f"[MSG] Processing: {message[:100]}...")
        
        # Get conversation state
        conversation_state = get_or_create_session(session_id)
        
        # Add user message to history
        conversation_state.add_message("user", message)
        
        print(f"Session {session_id} processing")
        
        # DETECT UPLOADED CONTENT: Check if message contains file content markers
        # If so, extract and store it for use in follow-up messages
        import re
        # Look for explicit file content markers
        file_content_match = re.search(r'=== FILE CONTENT ===\n(.*?)\n=== END FILE CONTENT ===', message, re.DOTALL)
        
        if file_content_match:
            uploaded_content = file_content_match.group(1).strip()
            if uploaded_content:  # Only store if non-empty
                conversation_state.set_uploaded_content(uploaded_content)
                print(f"[STORED] Uploaded content: {len(uploaded_content)} characters")
        
        # If user is asking for review/analysis and we have stored content, inject it
        processed_message = message
        if conversation_state.get_uploaded_content():
            # Check if message is asking for analysis/review but doesn't include the content
            review_keywords = ['review', 'analyze', 'check', 'examine', 'assess', 'compliance', 'full comp', 'it', 'this', 'that']
            has_review_keyword = any(kw in message.lower() for kw in review_keywords)
            no_file_content_marker = "=== FILE CONTENT ===" not in message
            
            if has_review_keyword and no_file_content_marker:
                # Inject the stored content
                stored_content = conversation_state.get_uploaded_content()
                processed_message = f"{message}\n\n=== FILE CONTENT (from previous upload) ===\n{stored_content}\n=== END FILE CONTENT ==="
                print(f"[INJECT] Injected stored content ({len(stored_content)} chars) for follow-up analysis")
        
        # SIMPLE APPROACH: Let Claude handle everything naturally
        # Like claude.com - just send the message and let Claude decide what to do
        # No classification, no trigger words, no routing logic
        
        # Always use the agent
        use_agent = True
        
        # Get recent chat history for context
        recent_history = [
            {"role": msg["role"], "content": msg["content"]} 
            for msg in conversation_state.get_recent_messages(limit=10) 
            if msg["content"].strip()
        ]
        
        print(f"[Processing] Sending to agent for natural handling")
        
        # Build context from conversation state and prepend to user input
        from agent_runtime import build_context_string
        context_str = build_context_string(conversation_state)
        
        # Prepend context to user input if context exists
        if context_str and context_str.strip():
            processed_message_with_context = f"{context_str}\n\nUser message: {processed_message}"
        else:
            processed_message_with_context = processed_message
        
        try:
            # Handle streaming vs non-streaming
            if streaming:
                # Use streaming with callbacks for real-time response
                from flask import Response, stream_with_context
                from langchain_core.callbacks.base import BaseCallbackHandler
                import queue
                import threading
                
                # Custom callback that puts tokens in a queue
                class StreamingCallbackHandler(BaseCallbackHandler):
                    def __init__(self, token_queue):
                        self.token_queue = token_queue
                    
                    def on_llm_new_token(self, token: str, **kwargs):
                        """Called when LLM generates a new token"""
                        self.token_queue.put(token)
                    
                    def on_llm_end(self, response, **kwargs):
                        """Called when LLM finishes"""
                        self.token_queue.put(None)  # Signal completion
                    
                    def on_llm_error(self, error: Exception, **kwargs):
                        """Called when LLM encounters an error"""
                        self.token_queue.put(f"ERROR:{str(error)}")
                        self.token_queue.put(None)
                
                def generate():
                    """Generator for Server-Sent Events"""
                    token_queue = queue.Queue()
                    full_response = ""
                    
                    # Create callback handler
                    callback = StreamingCallbackHandler(token_queue)
                    
                    # Run LLM in separate thread
                    def run_llm():
                        try:
                            agent_executor.invoke(
                                {
                                    "input_text": processed_message_with_context,
                                    "chat_history": recent_history
                                },
                                config={"callbacks": [callback]}
                            )
                        except Exception as e:
                            token_queue.put(f"ERROR:{str(e)}")
                            token_queue.put(None)
                    
                    llm_thread = threading.Thread(target=run_llm)
                    llm_thread.start()
                    
                    # Stream tokens as they arrive
                    try:
                        while True:
                            token = token_queue.get()
                            
                            if token is None:
                                # LLM finished
                                break
                            
                            if isinstance(token, str) and token.startswith("ERROR:"):
                                # Error occurred
                                error_msg = token[6:]
                                yield f"data: {json.dumps({'error': error_msg})}\n\n"
                                break
                            
                            # Send token to client
                            full_response += token
                            yield f"data: {json.dumps({'chunk': token})}\n\n"
                        
                        # Store analysis in memory
                        if any(keyword in processed_message.lower() for keyword in ['analyze', 'review', 'check', 'compliance', 'issues', 'find', 'problems', 'errors']):
                            issues = []
                            issue_matches = re.findall(r'[•\-\*]\s*\*\*(.+?)\*\*', full_response)
                            if issue_matches:
                                issues = issue_matches
                            else:
                                issue_matches = re.findall(r'\d+\.\s+(.+?)(?:\n|$)', full_response)
                                if issue_matches:
                                    issues = [match.split('\n')[0].strip() for match in issue_matches]
                            
                            if issues or 'compliance' in full_response.lower() or 'issue' in full_response.lower():
                                conversation_state.set_last_analysis(
                                    analysis_summary=full_response[:200],
                                    issues=issues
                                )
                                print(f"[MEMORY] Stored analysis with {len(issues)} issues")
                        
                        # Add to history
                        if full_response.strip():
                            conversation_state.add_message("assistant", full_response)
                        
                        # Extract citations
                        citations_dict = extract_and_parse_citations(full_response)
                        
                        # Send completion
                        yield f"data: {json.dumps({'done': True, 'full_response': full_response, 'citations': citations_dict})}\n\n"
                        
                    finally:
                        llm_thread.join(timeout=1)
                
                return Response(
                    stream_with_context(generate()),
                    mimetype='text/event-stream',
                    headers={
                        'Cache-Control': 'no-cache',
                        'X-Accel-Buffering': 'no',
                        'Connection': 'keep-alive'
                    }
                )
            
            else:
                # Non-streaming (traditional) response
                raw_response = agent_executor.invoke({
                    "input_text": processed_message_with_context,
                    "chat_history": recent_history
                })
                
                # Extract text from response (handles both dict and AIMessage objects)
                if isinstance(raw_response, dict):
                    response_text = raw_response.get('text', '').strip()
                elif hasattr(raw_response, 'content'):
                    # AIMessage object from LangChain
                    response_text = raw_response.content.strip()
                else:
                    response_text = str(raw_response).strip()
                
                if not response_text.strip():
                    response_text = "Hello! I'm EyeQ, your MLR (Marketing Legal Review) compliance agent. I help review marketing materials for regulatory compliance. You can upload documents or paste promotional content, and I'll perform a comprehensive analysis checking for claims, disclaimers, regulatory language, consistency, tone, and audience appropriateness. How can I help you today?"
                
                print(f"[OK] Agent response: {response_text[:100]}...")
                analysis_data = None
                
                # Store analysis in memory if it was a compliance review
                if any(keyword in processed_message.lower() for keyword in ['analyze', 'review', 'check', 'compliance', 'issues', 'find', 'problems', 'errors']):
                    import re
                    issues = []
                    # Try to extract from markdown headers or bullet points
                    issue_matches = re.findall(r'[•\-\*]\s*\*\*(.+?)\*\*', response_text)
                    if issue_matches:
                        issues = issue_matches
                    else:
                        # Fallback: try to extract from numbered items
                        issue_matches = re.findall(r'\d+\.\s+(.+?)(?:\n|$)', response_text)
                        if issue_matches:
                            issues = [match.split('\n')[0].strip() for match in issue_matches]
                    
                    # Store in conversation state
                    if issues or 'compliance' in response_text.lower() or 'issue' in response_text.lower():
                        conversation_state.set_last_analysis(
                            analysis_summary=response_text[:200],
                            issues=issues
                        )
                        print(f"[MEMORY] Stored analysis with {len(issues)} issues")
                
                # Add response to history (non-streaming only)
                if response_text.strip():
                    conversation_state.add_message("assistant", response_text)
            
        except Exception as agent_error:
            print(f"Agent error: {agent_error}")
            response_text = "Hello! I'm EyeQ, your MLR (Marketing Legal Review) compliance agent. I help review marketing materials for regulatory compliance. How can I help?"
            analysis_data = None
        
        # Typed response model
        class AnalyzeResponse(BaseModel):
            response: str
            analysis: dict | None
            full_response: str
            formatted: bool
            citations: dict

        # Extract citations from response
        citations_dict = extract_and_parse_citations(response_text)
        
        try:
            payload = AnalyzeResponse(
                response=response_text,
                analysis=analysis_data,
                full_response=response_text,
                formatted=True,
                citations=citations_dict
            )
            return jsonify(payload.model_dump())
        except ValidationError as ve:
            # Fallback to untyped JSON on validation errors
            return jsonify({
                "response": response_text,
                "analysis": analysis_data,
                "full_response": response_text,
                "formatted": True,
                "citations": citations_dict
            })
            
    except Exception as e:
        print(f"[ERROR] Analyze error: {e}")
        print(traceback.format_exc())
        return jsonify({
            "error": str(e),
            "response": "I apologize, but an error occurred. Please try again."
        }), 500

@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file uploads"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        # Generate unique filename
        file_id = str(uuid.uuid4())
        file_extension = os.path.splitext(file.filename)[1].lower()
        filename = f"{file_id}{file_extension}"
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        
        # Save file
        file.save(filepath)
        
        # Extract content based on file type
        content = ""
        ocr_method = None
        try:
            if file_extension == '.txt':
                with open(filepath, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif file_extension == '.docx':
                content = read_docx(filepath)
            elif file_extension == '.pdf':
                with pdfplumber.open(filepath) as pdf:
                    content = "\n".join([page.extract_text() or "" for page in pdf.pages])
            elif file_extension == '.pptx':
                prs = Presentation(filepath)
                content = "\n".join([
                    "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    for slide in prs.slides
                ])
            elif is_image_file(file.filename):
                # NEW: OCR for images (screenshots, promotional images, etc.)
                print(f"[IMG] Image detected: {file.filename}, performing OCR...")
                ocr_result = extract_text_from_image(filepath, method='auto', use_claude_fallback=True)
                content = ocr_result['text']
                ocr_method = ocr_result['method']
                
                # Add metadata about OCR process
                if ocr_result['confidence'] != 'failed':
                    content = f"[Image Text Extracted via {ocr_method.upper()}]\n\n{content}"
                    print(f"[OK] OCR successful using {ocr_method}: {len(content)} characters")
                else:
                    print(f"[FAIL] OCR failed for {file.filename}")
            else:
                content = f"[File: {file.filename} - Unsupported format]"
        except Exception as e:
            print(f"Error extracting content from {file.filename}: {e}")
            content = f"[File: {file.filename} - Content extraction failed]"
        
        # Clean up file
        try:
            os.remove(filepath)
        except:
            pass
        
        return jsonify({
            "content": content,
            "filename": file.filename,
            "file_id": file_id,
            "ocr_method": ocr_method,
            "is_image": is_image_file(file.filename)
        })
        
    except Exception as e:
        print(f"Error in upload: {str(e)}")
        return jsonify({"error": f"Upload error: {str(e)}"}), 500

@app.route('/api/compare', methods=['POST']) # uses comparison utilities and returns a markdown report + stats
def compare_documents():
    """Compare two documents for compliance differences"""
    try:
        data = request.get_json()
        doc1_content = data.get('document1', '')
        doc2_content = data.get('document2', '')
        doc1_name = data.get('doc1_name', 'Document 1')
        doc2_name = data.get('doc2_name', 'Document 2')
        
        if not doc1_content or not doc2_content:
            return jsonify({"error": "Both documents are required"}), 400
        
        print(f"Comparing documents: {doc1_name} vs {doc2_name}")
        
        # Calculate text similarity
        similarity = calculate_text_similarity(doc1_content, doc2_content)
        
        # Extract and compare claims
        claims1 = extract_claims_from_text(doc1_content)
        claims2 = extract_claims_from_text(doc2_content)
        claim_comparison = compare_claims(claims1, claims2)
        
        # Get compliance analysis for both documents (simplified - you can integrate full analysis)
        # For now, we'll create mock analysis structures
        analysis1 = parse_structured_response(doc1_content)
        analysis2 = parse_structured_response(doc2_content)
        
        # Analyze compliance differences
        compliance_comparison = analyze_compliance_differences(analysis1, analysis2)
        
        # Format comparison report
        comparison_report = format_comparison_report(
            doc1_name,
            doc2_name,
            similarity,
            claim_comparison,
            compliance_comparison
        )
        
        return jsonify({
            "comparison_report": comparison_report,
            "similarity": round(similarity * 100, 1),
            "claim_comparison": claim_comparison,
            "compliance_comparison": compliance_comparison
        })
        
    except Exception as e:
        print(f"Error in compare: {str(e)}")
        print(traceback.format_exc())
        return jsonify({"error": f"Comparison error: {str(e)}"}), 500

@app.route('/api/knowledge', methods=['GET']) # Search the KB and returns a formatted markdown snippet + raw results
def search_knowledge():
    """Search the regulatory knowledge base"""
    try:
        query = request.args.get('q', '')
        
        if not query:
            # Return overview of available topics
            return jsonify({
                "topics": list(REGULATORY_KNOWLEDGE.keys()),
                "scenarios": list(COMPLIANCE_SCENARIOS.keys()),
                "message": "Provide a query parameter 'q' to search the knowledge base"
            })
        
        results = search_knowledge_base(query)
        formatted_response = format_knowledge_base_response(query, results)
        
        return jsonify({
            "query": query,
            "results_count": len(results),
            "formatted_response": formatted_response,
            "results": results
        })
        
    except Exception as e:
        print(f"Error in knowledge search: {str(e)}")
        return jsonify({"error": f"Knowledge search error: {str(e)}"}), 500

@app.route('/api/share', methods=['POST']) 
def create_share():
    """Create a shareable link for a conversation"""
    try:
        data = request.get_json()
        conversation_id = data.get('conversation_id', '')
        messages = data.get('messages', [])
        analysis_summary = data.get('analysis_summary')
        
        if not conversation_id or not messages:
            return jsonify({"error": "Conversation ID and messages are required"}), 400
        
        share_info = create_shareable_link(conversation_id, messages, analysis_summary)
        
        return jsonify(share_info)
        
    except Exception as e:
        print(f"Error creating share: {str(e)}")
        return jsonify({"error": f"Share creation error: {str(e)}"}), 500

@app.route('/api/share/<share_id>', methods=['GET'])
def get_share(share_id):
    """Retrieve a shared analysis"""
    try:
        shared_analysis = get_shared_analysis(share_id)
        return jsonify(shared_analysis)
    except ValueError as e:
        return jsonify({"error": str(e)}), 404
    except Exception as e:
        return jsonify({"error": f"Error retrieving share: {str(e)}"}), 500

@app.route('/api/share/<share_id>/comment', methods=['POST'])
def add_share_comment(share_id):
    """Add a comment to a shared analysis"""
    try:
        data = request.get_json()
        user_name = data.get('user_name', 'Anonymous')
        comment_text = data.get('comment_text', '')
        message_index = data.get('message_index')
        
        if not comment_text:
            return jsonify({"error": "Comment text is required"}), 400
        
        comment = add_comment(share_id, user_name, comment_text, message_index)
        return jsonify(comment)
        
    except ValueError as e:
        return jsonify({"error": str(e)}), 404
    except Exception as e:
        return jsonify({"error": f"Error adding comment: {str(e)}"}), 500

@app.route('/api/export/<conversation_id>', methods=['POST'])
def export_conv(conversation_id):
    """Export a conversation in specified format"""
    try:
        data = request.get_json()
        messages = data.get('messages', [])
        export_format = data.get('format', 'json')
        
        if not messages:
            return jsonify({"error": "Messages are required"}), 400
        
        exported_content = export_conversation(conversation_id, messages, export_format)
        
        return jsonify({
            "conversation_id": conversation_id,
            "format": export_format,
            "content": exported_content
        })
        
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        return jsonify({"error": f"Export error: {str(e)}"}), 500

@app.route('/api/collaboration/stats/<conversation_id>', methods=['GET'])
def get_collab_stats(conversation_id):
    """Get collaboration statistics for a conversation"""
    try:
        stats = get_collaboration_stats(conversation_id)
        return jsonify(stats)
    except Exception as e:
        return jsonify({"error": f"Error getting stats: {str(e)}"}), 500

@app.route('/api/comprehensive-review', methods=['POST'])
def comprehensive_review():
    """
    Comprehensive MLR compliance analysis endpoint.
    Performs exhaustive line-by-line analysis across 6 categories.
    Returns structured table format with all compliance issues and recommendations.
    """
    try:
        data = request.get_json()
        material_text = data.get('material_text', '')
        product_name = data.get('product_name')
        session_id = data.get('session_id', 'default')
        streaming = data.get('streaming', True)  # Enabled by default for better UX
        
        if not material_text:
            return jsonify({"error": "Material text is required for comprehensive review"}), 400
        
        print(f"[REVIEW] Comprehensive review requested (session: {session_id})")
        print(f"[LEN] Material length: {len(material_text)} characters")
        
        # Get conversation state
        conversation_state = get_or_create_session(session_id)
        
        # Add user request to history
        conversation_state.add_message("user", f"Please perform comprehensive MLR review of this material: {material_text[:100]}...")
        
        # Get recent chat history
        recent_history = [
            {"role": msg["role"], "content": msg["content"]} 
            for msg in conversation_state.get_recent_messages(limit=10) 
            if msg["content"].strip()
        ]
        
        print(f"[AGENT] Using comprehensive analysis agent...")
        
        # Build context from conversation state and prepend to user input
        from agent_runtime import build_context_string
        context_str = build_context_string(conversation_state)
        
        # Prepend context to user input if context exists
        analysis_prompt = f"Perform comprehensive MLR compliance analysis on this promotional material:\n\n{material_text}"
        if context_str and context_str.strip():
            analysis_prompt_with_context = f"{context_str}\n\n{analysis_prompt}"
        else:
            analysis_prompt_with_context = analysis_prompt
        
        try:
            # Handle streaming vs non-streaming
            if streaming:
                # Use streaming with callbacks for real-time comprehensive analysis
                from flask import Response, stream_with_context
                from langchain_core.callbacks.base import BaseCallbackHandler
                import queue
                import threading
                
                # Custom callback that puts tokens in a queue
                class StreamingCallbackHandler(BaseCallbackHandler):
                    def __init__(self, token_queue):
                        self.token_queue = token_queue
                    
                    def on_llm_new_token(self, token: str, **kwargs):
                        """Called when LLM generates a new token"""
                        self.token_queue.put(token)
                    
                    def on_llm_end(self, response, **kwargs):
                        """Called when LLM finishes"""
                        self.token_queue.put(None)  # Signal completion
                
                def generate():
                    """Generator for Server-Sent Events"""
                    token_queue = queue.Queue()
                    full_response = ""
                    
                    # Create callback handler
                    callback = StreamingCallbackHandler(token_queue)
                    
                    # Run LLM in separate thread
                    def run_llm():
                        try:
                            comprehensive_agent_executor.invoke(
                                {
                                    "input_text": analysis_prompt_with_context,
                                    "chat_history": recent_history
                                },
                                config={"callbacks": [callback]}
                            )
                        except Exception as e:
                            token_queue.put(f"ERROR:{str(e)}")
                            token_queue.put(None)
                    
                    llm_thread = threading.Thread(target=run_llm)
                    llm_thread.start()
                    
                    # Stream tokens as they arrive
                    try:
                        while True:
                            token = token_queue.get()
                            
                            if token is None:
                                # LLM finished
                                break
                            
                            if isinstance(token, str) and token.startswith("ERROR:"):
                                # Error occurred
                                error_msg = token[6:]
                                yield f"data: {json.dumps({'error': error_msg})}\n\n"
                                break
                            
                            # Send token to client
                            full_response += token
                            yield f"data: {json.dumps({'chunk': token})}\n\n"
                        
                        print(f"[OK] Comprehensive analysis streaming complete ({len(full_response)} characters)")
                        
                        # Store comprehensive analysis in memory
                        issues = []
                        issue_matches = re.findall(r'[•\-\*]\s*\*\*(.+?)\*\*', full_response)
                        if issue_matches:
                            issues = issue_matches
                        else:
                            issue_matches = re.findall(r'\d+\.\s+(.+?)(?:\n|$)', full_response)
                            if issue_matches:
                                issues = [match.split('\n')[0].strip() for match in issue_matches]
                        
                        if issues or 'compliance' in full_response.lower():
                            conversation_state.set_last_analysis(
                                analysis_summary=full_response[:200],
                                issues=issues
                            )
                            print(f"[MEMORY] Stored comprehensive analysis with {len(issues)} issues")
                        
                        # Add to history
                        if full_response.strip():
                            conversation_state.add_message("assistant", full_response)
                        
                        # Send completion signal
                        yield f"data: {json.dumps({'done': True, 'full_response': full_response})}\n\n"
                        
                    finally:
                        llm_thread.join(timeout=1)
                
                return Response(
                    stream_with_context(generate()),
                    mimetype='text/event-stream',
                    headers={
                        'Cache-Control': 'no-cache',
                        'X-Accel-Buffering': 'no',
                        'Connection': 'keep-alive'
                    }
                )
            
            else:
                # Non-streaming (traditional) comprehensive analysis
                raw_response = comprehensive_agent_executor.invoke({
                    "input_text": analysis_prompt_with_context,
                    "chat_history": recent_history
                })
                
                # Extract response (handles both dict and AIMessage objects)
                full_text = ""
                if isinstance(raw_response, dict):
                    output = raw_response.get("output", raw_response.get("text", ""))
                    # Handle various response formats
                    if isinstance(output, list):
                        for item in output:
                            if isinstance(item, dict) and 'text' in item:
                                full_text += item['text'] + "\n"
                            elif isinstance(item, str):
                                full_text += item + "\n"
                    elif isinstance(output, dict) and 'text' in output:
                        full_text = output['text']
                    elif isinstance(output, str):
                        full_text = output
                    else:
                        full_text = str(output)
                elif hasattr(raw_response, 'content'):
                    # AIMessage object from LangChain
                    full_text = raw_response.content
                else:
                    full_text = str(raw_response)
                
                full_text = full_text.strip()
                print(f"[OK] Comprehensive analysis complete ({len(full_text)} characters)")
                
                # Store comprehensive analysis in memory
                import re
                issues = []
                # Extract issues from comprehensive analysis response
                issue_matches = re.findall(r'[•\-\*]\s*\*\*(.+?)\*\*', full_text)
                if issue_matches:
                    issues = issue_matches
                else:
                    # Fallback: try to extract from numbered items
                    issue_matches = re.findall(r'\d+\.\s+(.+?)(?:\n|$)', full_text)
                    if issue_matches:
                        issues = [match.split('\n')[0].strip() for match in issue_matches]
                
                # Store comprehensive analysis results (non-streaming only)
                conversation_state.set_last_analysis(
                    analysis_summary=full_text[:200],
                    issues=issues
                )
                print(f"[MEMORY] Stored comprehensive analysis with {len(issues)} issues")
                
                # Add response to history (non-streaming only)
                if full_text.strip():
                    conversation_state.add_message("assistant", full_text)
            
            # Parse structured data if possible
            try:
                # Try to extract JSON if embedded in response
                json_start = full_text.find('{')
                json_end = full_text.rfind('}') + 1
                if json_start != -1 and json_end > json_start:
                    json_str = full_text[json_start:json_end]
                    structured_data = json.loads(json_str)
                else:
                    structured_data = None
            except:
                structured_data = None
            
            # Build comprehensive response
            response_payload = {
                "status": "success",
                "comprehensive_analysis": full_text,
                "structured_data": structured_data,
                "session_id": session_id,
                "analysis_type": "comprehensive_mlr_review",
                "timestamp": datetime.now(timezone.utc).isoformat()
            }
            
            return jsonify(response_payload)
            
        except Exception as agent_error:
            print(f"[FAIL] Comprehensive agent execution failed: {agent_error}")
            print(traceback.format_exc())
            return jsonify({
                "error": "Comprehensive analysis failed",
                "details": str(agent_error)
            }), 500
        
    except Exception as e:
        print(f"[FAIL] Error in comprehensive_review: {str(e)}")
        print(traceback.format_exc())
        return jsonify({"error": f"Comprehensive review error: {str(e)}"}), 500

# Feedback storage (in-memory for now, can be persisted to file/database later)
feedback_storage = {}

@app.route('/api/feedback', methods=['POST'])
def submit_feedback():
    """Submit user feedback (like/dislike) for agent responses"""
    try:
        data = request.get_json()
        message_id = data.get('message_id')
        message_content = data.get('message_content', '')
        feedback_type = data.get('feedback_type')  # 'liked' or 'disliked'
        conversation_id = data.get('conversation_id')
        
        if not message_id or not feedback_type:
            return jsonify({"error": "message_id and feedback_type are required"}), 400
        
        if feedback_type not in ['liked', 'disliked']:
            return jsonify({"error": "feedback_type must be 'liked' or 'disliked'"}), 400
        
        # Store feedback
        feedback_entry = {
            'message_id': message_id,
            'message_content': message_content[:500],  # Store first 500 chars for context
            'feedback_type': feedback_type,
            'conversation_id': conversation_id,
            'timestamp': datetime.now(timezone.utc).isoformat()
        }
        
        # Store in memory (keyed by message_id)
        feedback_storage[message_id] = feedback_entry
        
        # Also store by conversation for learning
        if conversation_id:
            if conversation_id not in feedback_storage:
                feedback_storage[conversation_id] = []
            if isinstance(feedback_storage[conversation_id], list):
                feedback_storage[conversation_id].append(feedback_entry)
        
        # Save feedback to file for persistence (optional)
        feedback_file = 'feedback_log.json'
        try:
            import json as json_lib
            # Load existing feedback
            existing_feedback = []
            if os.path.exists(feedback_file):
                with open(feedback_file, 'r') as f:
                    existing_feedback = json_lib.load(f)
            
            # Append new feedback
            existing_feedback.append(feedback_entry)
            
            # Save back to file
            with open(feedback_file, 'w') as f:
                json_lib.dump(existing_feedback, f, indent=2)
        except Exception as e:
            print(f"Warning: Could not save feedback to file: {e}")
        
        print(f"[FEEDBACK] {feedback_type.upper()} for message {message_id}")
        
        return jsonify({
            "status": "success",
            "message": f"Feedback recorded: {feedback_type}",
            "feedback_id": message_id
        })
        
    except Exception as e:
        print(f"Error in feedback endpoint: {str(e)}")
        print(traceback.format_exc())
        return jsonify({"error": f"Feedback error: {str(e)}"}), 500

@app.route('/api/feedback/learn', methods=['GET'])
def get_feedback_for_learning():
    """Get feedback data for agent learning (can be called by agent)"""
    try:
        # Return recent feedback
        feedback_file = 'feedback_log.json'
        if os.path.exists(feedback_file):
            import json as json_lib
            with open(feedback_file, 'r') as f:
                feedback_data = json_lib.load(f)
                # Return last 50 feedback entries
                return jsonify({
                    "feedback": feedback_data[-50:] if len(feedback_data) > 50 else feedback_data,
                    "total_count": len(feedback_data)
                })
        else:
            return jsonify({"feedback": [], "total_count": 0})
    except Exception as e:
        return jsonify({"error": f"Error loading feedback: {str(e)}"}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    """Check system health"""
    try:
        return jsonify({
            "status": "healthy",
            "agent": "EyeQ (Working Version)",
            "database": "localStorage",
            "timestamp": datetime.now(timezone.utc).isoformat()
        })
    except Exception as e:
        return jsonify({
            "status": "unhealthy",
            "agent": "EyeQ (Working Version)",
            "database": "localStorage",
            "error": str(e),
            "timestamp": datetime.now(timezone.utc).isoformat()
        }), 500

if __name__ == '__main__':
    print("=" * 60)
    print(" EyeQ - Web Interface (WORKING VERSION)")
    print("=" * 60)
    print("Make sure your .env file contains:")
    print("- ANTHROPIC_API_KEY=your_api_key")
    print()
    print("[OK] Using EXACT same agent system as main.py!")
    print("[AGENT] AI Analysis: Claude 3.5 Sonnet with compliance tools")
    print("[DATA] Data persistence: localStorage (frontend)")
    print(f"[START] Starting server on http://127.0.0.1:5000")
    print("=" * 60)

    app.run(debug=True, host='127.0.0.1', port=5000)
